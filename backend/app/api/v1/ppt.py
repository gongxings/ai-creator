"""
PPT生成API路由
"""
from typing import List
from fastapi import APIRouter, Depends, HTTPException, BackgroundTasks
from sqlalchemy.orm import Session

from app.core.database import get_db
from app.core.security import get_current_user
from app.models.user import User
from app.models.creation import Creation
from app.schemas.creation import CreationCreate, CreationResponse
from app.services.ai.factory import AIServiceFactory

router = APIRouter()


@router.post("/generate", response_model=CreationResponse)
async def generate_ppt(
    request: CreationCreate,
    background_tasks: BackgroundTasks,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """
    生成PPT
    
    支持的生成方式：
    - 主题生成：提供主题和大纲
    - 文档转换：上传文档内容
    """
    try:
        # 创建创作记录
        creation = Creation(
            user_id=current_user.id,
            type="ppt",
            title=request.title or "未命名PPT",
            input_data=request.input_data,
            status="processing"
        )
        db.add(creation)
        db.commit()
        db.refresh(creation)
        
        # 异步生成PPT
        background_tasks.add_task(
            _generate_ppt_task,
            creation.id,
            request.input_data,
            request.config
        )
        
        return creation
        
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"生成PPT失败: {str(e)}")


async def _generate_ppt_task(creation_id: int, input_data: dict, config: dict = None):
    """
    异步生成PPT任务
    """
    from app.core.database import SessionLocal
    
    db = SessionLocal()
    try:
        creation = db.query(Creation).filter(Creation.id == creation_id).first()
        if not creation:
            return
        
        # 获取AI服务
        ai_service = AIServiceFactory.get_service(config.get("model", "openai") if config else "openai")
        
        # 构建提示词
        prompt = _build_ppt_prompt(input_data)
        
        # 调用AI生成PPT大纲和内容
        result = await ai_service.generate(prompt)
        
        # 解析生成结果
        ppt_data = _parse_ppt_result(result)
        
        # 更新创作记录
        creation.content = ppt_data
        creation.status = "completed"
        db.commit()
        
    except Exception as e:
        if creation:
            creation.status = "failed"
            creation.error_message = str(e)
            db.commit()
    finally:
        db.close()


def _build_ppt_prompt(input_data: dict) -> str:
    """
    构建PPT生成提示词
    """
    generation_type = input_data.get("type", "theme")
    
    if generation_type == "theme":
        # 主题生成PPT
        theme = input_data.get("theme", "")
        outline = input_data.get("outline", "")
        slides_count = input_data.get("slides_count", 10)
        
        prompt = f"""请根据以下主题和大纲生成一个专业的PPT内容：

主题：{theme}
大纲：{outline}
幻灯片数量：{slides_count}

要求：
1. 生成完整的PPT结构，包括标题页、目录页、内容页和结束页
2. 每页包含标题、要点和详细说明
3. 内容要专业、清晰、有逻辑性
4. 适当使用数据、案例和图表说明
5. 返回JSON格式，包含slides数组，每个slide包含：
   - title: 标题
   - content: 内容要点（数组）
   - notes: 演讲备注
   - layout: 布局类型（title/content/two-column/image等）

请直接返回JSON格式的PPT内容。"""

    elif generation_type == "document":
        # 文档转PPT
        document = input_data.get("document", "")
        slides_count = input_data.get("slides_count", 10)
        
        prompt = f"""请将以下文档内容转换为PPT格式：

文档内容：
{document}

要求：
1. 提取文档的核心内容和结构
2. 生成约{slides_count}页幻灯片
3. 每页包含标题、要点和详细说明
4. 保持原文档的逻辑结构
5. 返回JSON格式，包含slides数组

请直接返回JSON格式的PPT内容。"""

    else:
        # 大纲生成PPT
        outline = input_data.get("outline", "")
        
        prompt = f"""请根据以下大纲生成PPT内容：

大纲：
{outline}

要求：
1. 根据大纲结构生成幻灯片
2. 每个大纲要点对应一页或多页幻灯片
3. 内容要详细、专业
4. 返回JSON格式

请直接返回JSON格式的PPT内容。"""
    
    return prompt


def _parse_ppt_result(result: str) -> dict:
    """
    解析PPT生成结果
    """
    import json
    import re
    
    try:
        # 尝试直接解析JSON
        return json.loads(result)
    except:
        # 如果不是纯JSON，尝试提取JSON部分
        json_match = re.search(r'\{[\s\S]*\}', result)
        if json_match:
            return json.loads(json_match.group())
        
        # 如果都失败，返回原始文本
        return {
            "slides": [
                {
                    "title": "生成的PPT内容",
                    "content": [result],
                    "notes": "",
                    "layout": "content"
                }
            ]
        }


@router.get("/{ppt_id}/download")
async def download_ppt(
    ppt_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """
    下载PPT文件
    """
    creation = db.query(Creation).filter(
        Creation.id == ppt_id,
        Creation.user_id == current_user.id,
        Creation.type == "ppt"
    ).first()
    
    if not creation:
        raise HTTPException(status_code=404, detail="PPT不存在")
    
    if creation.status != "completed":
        raise HTTPException(status_code=400, detail="PPT尚未生成完成")
    
    # 生成PPT文件
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        import os
        from datetime import datetime
        
        # 创建PPT对象
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # 获取PPT数据
        ppt_data = creation.content
        slides_data = ppt_data.get("slides", [])
        
        # 生成每一页幻灯片
        for slide_data in slides_data:
            layout_type = slide_data.get("layout", "content")
            
            # 根据布局类型选择幻灯片布局
            if layout_type == "title":
                slide_layout = prs.slide_layouts[0]  # 标题页
            elif layout_type == "two-column":
                slide_layout = prs.slide_layouts[3]  # 两栏布局
            else:
                slide_layout = prs.slide_layouts[1]  # 标题和内容
            
            slide = prs.slides.add_slide(slide_layout)
            
            # 设置标题
            title = slide.shapes.title
            title.text = slide_data.get("title", "")
            
            # 设置内容
            if layout_type != "title" and len(slide.shapes) > 1:
                content_shape = slide.shapes[1]
                if hasattr(content_shape, "text_frame"):
                    text_frame = content_shape.text_frame
                    text_frame.clear()
                    
                    content_items = slide_data.get("content", [])
                    for item in content_items:
                        p = text_frame.add_paragraph()
                        p.text = item
                        p.level = 0
        
        # 保存PPT文件
        upload_dir = "app/uploads/ppt"
        os.makedirs(upload_dir, exist_ok=True)
        
        filename = f"ppt_{ppt_id}_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx"
        filepath = os.path.join(upload_dir, filename)
        prs.save(filepath)
        
        return {
            "code": 200,
            "message": "success",
            "data": {
                "download_url": f"/uploads/ppt/{filename}",
                "filename": filename,
                "ppt_data": creation.content
            }
        }
        
    except ImportError:
        # 如果没有安装python-pptx，返回JSON数据
        return {
            "code": 200,
            "message": "PPT数据已生成，但需要安装python-pptx库才能导出文件",
            "data": {
                "ppt_data": creation.content,
                "note": "请运行: pip install python-pptx"
            }
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"生成PPT文件失败: {str(e)}")
