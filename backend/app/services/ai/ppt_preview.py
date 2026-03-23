"""
PPT预览服务
生成PPT缩略图和预览
"""
import os
import logging
import tempfile
from typing import List, Optional
from pathlib import Path

logger = logging.getLogger(__name__)


class PPTPreviewService:
    """PPT预览服务"""
    
    def __init__(self):
        """初始化预览服务"""
        self.available = False
        try:
            # 尝试导入python-pptx
            from pptx import Presentation
            self.Presentation = Presentation
            self.available = True
        except ImportError:
            logger.warning("python-pptx not installed, preview generation will be unavailable")
    
    def generate_thumbnails(
        self,
        pptx_path: str,
        output_dir: str,
        width: int = 320,
        height: int = 180
    ) -> List[str]:
        """
        生成PPT缩略图
        
        注意：python-pptx本身不支持直接渲染PPT为图片。
        这里提供一个基础实现，实际部署时可能需要配合其他工具（如LibreOffice）。
        
        Args:
            pptx_path: PPTX文件路径
            output_dir: 输出目录
            width: 缩略图宽度
            height: 缩略图高度
            
        Returns:
            缩略图路径列表
        """
        if not self.available:
            raise RuntimeError("python-pptx library not available")
        
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"PPTX file not found: {pptx_path}")
        
        # 确保输出目录存在
        os.makedirs(output_dir, exist_ok=True)
        
        try:
            prs = self.Presentation(pptx_path)
            thumbnail_paths = []
            
            # 生成占位缩略图（实际项目中应使用LibreOffice等工具渲染）
            for idx, slide in enumerate(prs.slides):
                thumbnail_path = os.path.join(output_dir, f"slide_{idx + 1}.png")
                
                # 这里只是创建一个占位文件
                # 实际实现需要使用其他工具来渲染PPT为图片
                self._create_placeholder_thumbnail(thumbnail_path, idx + 1, len(prs.slides))
                
                thumbnail_paths.append(thumbnail_path)
            
            return thumbnail_paths
            
        except Exception as e:
            logger.error(f"Failed to generate thumbnails: {e}")
            raise
    
    def _create_placeholder_thumbnail(self, path: str, slide_num: int, total_slides: int):
        """
        创建占位缩略图
        
        Args:
            path: 输出路径
            slide_num: 当前幻灯片编号
            total_slides: 总幻灯片数
        """
        try:
            from PIL import Image, ImageDraw, ImageFont
            
            # 创建空白图片
            img = Image.new('RGB', (320, 180), color='#f0f0f0')
            draw = ImageDraw.Draw(img)
            
            # 绘制边框
            draw.rectangle([0, 0, 319, 179], outline='#cccccc', width=2)
            
            # 绘制文字
            text = f"Slide {slide_num}/{total_slides}"
            draw.text((140, 80), text, fill='#666666')
            
            img.save(path)
            
        except ImportError:
            # 如果没有Pillow，创建一个空文件
            with open(path, 'wb') as f:
                # 写入一个最小的PNG文件头
                f.write(b'\x89PNG\r\n\x1a\n')
    
    def extract_slide_text(self, pptx_path: str) -> List[dict]:
        """
        提取PPT中每页的文本内容
        
        Args:
            pptx_path: PPTX文件路径
            
        Returns:
            每页文本内容列表
        """
        if not self.available:
            raise RuntimeError("python-pptx library not available")
        
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"PPTX file not found: {pptx_path}")
        
        try:
            prs = self.Presentation(pptx_path)
            slides_text = []
            
            for idx, slide in enumerate(prs.slides):
                slide_text = {
                    "index": idx,
                    "title": "",
                    "content": []
                }
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if shape.shape_type == 13:  # TEXT_BOX
                            slide_text["content"].append(shape.text)
                        elif hasattr(shape, "placeholder_format"):
                            ph_type = shape.placeholder_format.type
                            if ph_type == 0:  # TITLE
                                slide_text["title"] = shape.text
                            else:
                                slide_text["content"].append(shape.text)
                
                slides_text.append(slide_text)
            
            return slides_text
            
        except Exception as e:
            logger.error(f"Failed to extract slide text: {e}")
            raise


# 全局实例
ppt_preview_service = PPTPreviewService()
