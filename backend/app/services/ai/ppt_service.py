"""
PPT生成Service - 基于Cookie的网页版实现
支持主题生成、大纲生成、文档转PPT等功能
"""
import httpx
import json
import logging
from typing import Dict, Any, Optional, List

from app.services.ai.cookie_based_service import CookieBasedAIService

logger = logging.getLogger(__name__)


class PPTGenerationService(CookieBasedAIService):
    """PPT生成基础服务"""
    
    def get_platform_name(self) -> str:
        return "ppt_generation"
    
    def get_check_url(self) -> str:
        return "https://www.doubao.com/"
    
    async def generate_ppt_outline(
        self,
        title: str,
        content: str,
        num_slides: int = 10,
        **kwargs
    ) -> Dict[str, Any]:
        """
        生成PPT大纲
        """
        # 需要实际的PPT生成API或使用python-pptx库
        return {
            "slides": [],
            "error": "PPT生成需要专用库支持"
        }
    
    async def generate_text(self, prompt: str, **kwargs) -> str:
        """不支持"""
        raise NotImplementedError("PPTGenerationService 不支持文本生成")
    
    async def generate_image(self, prompt: str, **kwargs) -> Dict[str, Any]:
        """不支持"""
        raise NotImplementedError("PPTGenerationService 不支持图片生成")
    
    async def generate_video(self, prompt: str, **kwargs) -> Dict[str, Any]:
        """不支持"""
        raise NotImplementedError("PPTGenerationService 不支持视频生成")


class DoubiaoPPTService(CookieBasedAIService):
    """豆包PPT生成服务"""
    
    BASE_URL = "https://www.doubao.com"
    DEFAULT_BOT_ID = "7358044466096914465"
    
    def get_platform_name(self) -> str:
        return "doubao_ppt"
    
    def get_check_url(self) -> str:
        return self.BASE_URL
    
    async def generate_ppt_outline(
        self,
        title: str,
        content: str,
        num_slides: int = 10,
        **kwargs
    ) -> Dict[str, Any]:
        """
        通过豆包生成PPT大纲
        """
        headers = self.get_headers()
        
        # 通过Chat API生成大纲
        prompt = f"""请帮我生成一份{num_slides}页的PPT大纲：

标题：{title}
内容描述：{content}

请按照以下格式生成每页的标题和要点：
第1页：标题页
- 主标题：...
- 副标题：...

第2页：概述
- 要点1：...
- 要点2：...
...

注意：
1. 每页包含标题和3-5个要点
2. 逻辑清晰，循序渐进
3. 最后一页为总结和展望
"""
        
        payload = {
            "user_input": prompt,
            "bot_id": self.DEFAULT_BOT_ID,
            "stream": False,
        }
        
        logger.info(f"Doubao PPT outline generation - title: {title}")
        
        try:
            async with httpx.AsyncClient(follow_redirects=True, timeout=120.0) as client:
                response = await client.post(
                    f"{self.BASE_URL}/api/chat/completions",
                    headers=headers,
                    json=payload,
                )
                
                if response.status_code == 401:
                    return {
                        "outline": None,
                        "error": "Cookie已过期，请重新登录授权"
                    }
                
                response.raise_for_status()
                data = response.json()
                
                # 提取PPT大纲
                outline = None
                if "choices" in data and len(data["choices"]) > 0:
                    choice = data["choices"][0]
                    if "message" in choice and "content" in choice["message"]:
                        outline = choice["message"]["content"]
                
                return {
                    "outline": outline,
                    "title": title,
                    "num_slides": num_slides,
                    "status": "outline_generated"
                }
                
        except Exception as e:
            logger.error(f"Doubao PPT generation failed: {e}")
            return {
                "outline": None,
                "error": str(e)
            }
    
    async def generate_text(self, prompt: str, **kwargs) -> str:
        """不支持"""
        raise NotImplementedError("DoubiaoPPTService 不支持文本生成")
    
    async def generate_image(self, prompt: str, **kwargs) -> Dict[str, Any]:
        """不支持"""
        raise NotImplementedError("DoubiaoPPTService 不支持图片生成")
    
    async def generate_video(self, prompt: str, **kwargs) -> Dict[str, Any]:
        """不支持"""
        raise NotImplementedError("DoubiaoPPTService 不支持视频生成")


class LocalPPTService:
    """本地PPT生成服务 - 使用python-pptx"""
    
    def __init__(self):
        """初始化本地PPT服务"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            self.Presentation = Presentation
            self.Inches = Inches
            self.Pt = Pt
        except ImportError:
            logger.warning("python-pptx not installed, PPT generation will be limited")
            self.available = False
        else:
            self.available = True
    
    def create_pptx_from_outline(
        self,
        title: str,
        outline: str,
        output_path: str = "/tmp/presentation.pptx"
    ) -> str:
        """
        从大纲创建PPT文件
        
        Args:
            title: PPT标题
            outline: PPT大纲（文本格式）
            output_path: 输出文件路径
            
        Returns:
            输出文件路径
        """
        if not self.available:
            raise RuntimeError("python-pptx library not available")
        
        # 创建演示文稿
        prs = self.Presentation()
        
        # 添加标题页
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title
        subtitle.text = "AI生成的演示文稿"
        
        # 解析大纲并添加页面
        # 这里需要实现大纲解析逻辑
        # 简化处理：按行分割并创建页面
        lines = outline.split("\n")
        current_slide_data = {"title": "", "content": []}
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # 简单的解析逻辑
            if line.startswith("第") and "页" in line:
                # 新页面
                if current_slide_data["title"]:
                    self._add_slide(prs, current_slide_data)
                current_slide_data = {"title": line, "content": []}
            elif line.startswith("-") or line.startswith("•"):
                # 要点
                current_slide_data["content"].append(line[1:].strip())
            elif line and not current_slide_data["title"]:
                current_slide_data["title"] = line
        
        # 添加最后一页
        if current_slide_data["title"]:
            self._add_slide(prs, current_slide_data)
        
        # 保存
        prs.save(output_path)
        logger.info(f"PPT created at: {output_path}")
        
        return output_path
    
    def _add_slide(self, prs, slide_data):
        """添加页面到演示文稿"""
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        title = slide.shapes.title
        title.text = slide_data["title"]
        
        # 添加要点
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        for idx, point in enumerate(slide_data["content"]):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.level = 0


# 工厂函数
def create_ppt_service(cookies: Dict[str, str], user_agent: Optional[str] = None) -> DoubiaoPPTService:
    """
    创建豆包PPT生成服务实例
    
    Args:
        cookies: Cookie字典
        user_agent: User-Agent（可选）
        
    Returns:
        DoubiaoPPTService实例
    """
    return DoubiaoPPTService(cookies=cookies, user_agent=user_agent)


def create_local_ppt_service() -> LocalPPTService:
    """
    创建本地PPT生成服务实例
    
    Returns:
        LocalPPTService实例
    """
    return LocalPPTService()
