"""
PPT模板解析服务
解析PPTX文件，提取布局元数据
"""
import os
import logging
from typing import Dict, Any, List, Optional
from pathlib import Path

logger = logging.getLogger(__name__)


class TemplateParser:
    """PPT模板解析器"""
    
    # 占位符类型映射
    PLACEHOLDER_TYPE_MAP = {
        0: "title",           # TITLE
        1: "body",            # BODY
        2: "subtitle",        # CENTER_TITLE
        3: "date",            # DATE
        4: "slide_number",    # SLIDE_NUMBER
        5: "footer",          # FOOTER
        6: "header",          # HEADER
        7: "object",          # OBJECT
        8: "chart",           # CHART
        9: "table",           # TABLE
        10: "clip_art",       # CLIP_ART
        11: "diagram",        # DIAGRAM
        12: "media",          # MEDIA
        13: "picture",        # PICTURE
    }
    
    def __init__(self):
        """初始化解析器"""
        self.available = False
        try:
            from pptx import Presentation
            from pptx.util import Emu
            self.Presentation = Presentation
            self.Emu = Emu
            self.available = True
        except ImportError:
            logger.warning("python-pptx not installed, template parsing will be unavailable")
    
    def parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """
        解析PPTX文件，提取布局元数据
        
        Args:
            file_path: PPTX文件路径
            
        Returns:
            布局元数据字典
        """
        if not self.available:
            raise RuntimeError("python-pptx library not available")
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PPTX file not found: {file_path}")
        
        try:
            prs = self.Presentation(file_path)
            
            slides_metadata = []
            
            for idx, slide in enumerate(prs.slides):
                slide_info = self._parse_slide(slide, idx)
                slides_metadata.append(slide_info)
            
            return {
                "slide_count": len(prs.slides),
                "slides": slides_metadata
            }
            
        except Exception as e:
            logger.error(f"Failed to parse PPTX file: {e}")
            raise
    
    def _parse_slide(self, slide, index: int) -> Dict[str, Any]:
        """
        解析单个幻灯片
        
        Args:
            slide: 幻灯片对象
            index: 幻灯片索引
            
        Returns:
            幻灯片信息字典
        """
        placeholders = []
        
        for placeholder in slide.placeholders:
            ph_info = self._parse_placeholder(placeholder)
            placeholders.append(ph_info)
        
        # 推断幻灯片类型
        slide_type = self._infer_slide_type(placeholders, index)
        
        # 获取布局名称
        layout_name = ""
        try:
            layout_name = slide.slide_layout.name
        except:
            pass
        
        return {
            "index": index,
            "type": slide_type,
            "layout_name": layout_name,
            "placeholders": placeholders
        }
    
    def _parse_placeholder(self, placeholder) -> Dict[str, Any]:
        """
        解析占位符
        
        Args:
            placeholder: 占位符对象
            
        Returns:
            占位符信息字典
        """
        ph_type = "unknown"
        try:
            ph_type = self.PLACEHOLDER_TYPE_MAP.get(placeholder.placeholder_format.type, "unknown")
        except:
            pass
        
        # 获取位置和大小
        bounds = None
        try:
            bounds = {
                "x": placeholder.left,
                "y": placeholder.top,
                "width": placeholder.width,
                "height": placeholder.height
            }
        except:
            pass
        
        # 获取标签
        label = ""
        try:
            label = placeholder.placeholder_format.name or ""
        except:
            pass
        
        # 检查是否支持项目符号
        supports_bullets = ph_type in ["body", "object"]
        
        return {
            "idx": placeholder.placeholder_format.idx,
            "type": ph_type,
            "label": label,
            "bounds": bounds,
            "supports_bullets": supports_bullets
        }
    
    def _infer_slide_type(self, placeholders: List[Dict], index: int) -> str:
        """
        根据占位符推断幻灯片类型
        
        Args:
            placeholders: 占位符列表
            index: 幻灯片索引
            
        Returns:
            幻灯片类型
        """
        if index == 0:
            # 第一页通常是标题页
            has_subtitle = any(p["type"] == "subtitle" for p in placeholders)
            if has_subtitle:
                return "title"
        
        # 检查是否有标题和正文
        has_title = any(p["type"] == "title" for p in placeholders)
        has_body = any(p["type"] == "body" for p in placeholders)
        
        if has_title and has_body:
            return "content"
        elif has_title and not has_body:
            return "section"
        
        return "content"
    
    def fill_template(
        self,
        template_path: str,
        outline: Dict[str, Any],
        output_path: str
    ) -> str:
        """
        使用大纲内容填充模板
        
        Args:
            template_path: 模板文件路径
            outline: 大纲内容
            output_path: 输出文件路径
            
        Returns:
            输出文件路径
        """
        if not self.available:
            raise RuntimeError("python-pptx library not available")
        
        try:
            prs = self.Presentation(template_path)
            
            slides = outline.get("slides", [])
            
            for idx, slide_content in enumerate(slides):
                if idx >= len(prs.slides):
                    break
                
                slide = prs.slides[idx]
                self._fill_slide(slide, slide_content)
            
            # 保存文件
            prs.save(output_path)
            
            return output_path
            
        except Exception as e:
            logger.error(f"Failed to fill template: {e}")
            raise
    
    def _fill_slide(self, slide, content: Dict[str, Any]):
        """
        填充单个幻灯片
        
        Args:
            slide: 幻灯片对象
            content: 内容字典
        """
        for placeholder in slide.placeholders:
            try:
                ph_type = self.PLACEHOLDER_TYPE_MAP.get(placeholder.placeholder_format.type, "unknown")
                idx = placeholder.placeholder_format.idx
                
                if ph_type == "title":
                    placeholder.text = content.get("title", "")
                elif ph_type == "subtitle":
                    placeholder.text = content.get("subtitle", "")
                elif ph_type == "body":
                    bullets = content.get("bullets", [])
                    if bullets:
                        self._fill_body_with_bullets(placeholder, bullets)
            except Exception as e:
                logger.warning(f"Failed to fill placeholder: {e}")
    
    def _fill_body_with_bullets(self, placeholder, bullets: List[str]):
        """
        使用项目符号填充正文占位符
        
        Args:
            placeholder: 占位符对象
            bullets: 项目符号列表
        """
        try:
            text_frame = placeholder.text_frame
            text_frame.clear()
            
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = bullet
                p.level = 0
        except Exception as e:
            logger.warning(f"Failed to fill bullets: {e}")


# 全局实例
template_parser = TemplateParser()
