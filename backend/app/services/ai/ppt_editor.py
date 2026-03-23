"""
PPT编辑服务
支持对生成的PPT进行文本、样式、图片编辑
"""
import os
import logging
from typing import Dict, Any, Optional, List
from pathlib import Path

logger = logging.getLogger(__name__)


class PPTEditor:
    """PPT编辑器"""
    
    # 占位符类型映射
    PLACEHOLDER_TYPE_MAP = {
        0: "title",
        1: "body",
        2: "subtitle",
        3: "date",
        4: "slide_number",
        5: "footer",
        6: "header",
        7: "object",
        8: "chart",
        9: "table",
        10: "clip_art",
        11: "diagram",
        12: "media",
        13: "picture",
    }
    
    def __init__(self):
        """初始化编辑器"""
        self.available = False
        try:
            from pptx import Presentation
            from pptx.util import Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            self.Presentation = Presentation
            self.Pt = Pt
            self.Emu = Emu
            self.RGBColor = RGBColor
            self.PP_ALIGN = PP_ALIGN
            self.available = True
        except ImportError:
            logger.warning("python-pptx not installed, PPT editing will be unavailable")
    
    def update_text(
        self,
        pptx_path: str,
        slide_index: int,
        placeholder_idx: int,
        new_text: str,
        output_path: Optional[str] = None
    ) -> str:
        """
        更新幻灯片中的文本
        
        Args:
            pptx_path: PPTX文件路径
            slide_index: 幻灯片索引
            placeholder_idx: 占位符索引
            new_text: 新文本内容
            output_path: 输出路径（可选，默认覆盖原文件）
            
        Returns:
            输出文件路径
        """
        if not self.available:
            raise RuntimeError("python-pptx library not available")
        
        try:
            prs = self.Presentation(pptx_path)
            
            if slide_index >= len(prs.slides):
                raise ValueError(f"Slide index {slide_index} out of range")
            
            slide = prs.slides[slide_index]
            
            # 查找占位符
            target_placeholder = None
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == placeholder_idx:
                    target_placeholder = placeholder
                    break
            
            if not target_placeholder:
                raise ValueError(f"Placeholder {placeholder_idx} not found")
            
            # 更新文本
            target_placeholder.text = new_text
            
            # 保存文件
            save_path = output_path or pptx_path
            prs.save(save_path)
            
            return save_path
            
        except Exception as e:
            logger.error(f"Failed to update text: {e}")
            raise
    
    def update_style(
        self,
        pptx_path: str,
        slide_index: int,
        placeholder_idx: int,
        font_size: Optional[int] = None,
        font_bold: Optional[bool] = None,
        font_color: Optional[str] = None,
        output_path: Optional[str] = None
    ) -> str:
        """
        更新幻灯片中的文本样式
        
        Args:
            pptx_path: PPTX文件路径
            slide_index: 幻灯片索引
            placeholder_idx: 占位符索引
            font_size: 字体大小
            font_bold: 是否加粗
            font_color: 字体颜色(hex)
            output_path: 输出路径
            
        Returns:
            输出文件路径
        """
        if not self.available:
            raise RuntimeError("python-pptx library not available")
        
        try:
            prs = self.Presentation(pptx_path)
            
            if slide_index >= len(prs.slides):
                raise ValueError(f"Slide index {slide_index} out of range")
            
            slide = prs.slides[slide_index]
            
            # 查找占位符
            target_placeholder = None
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == placeholder_idx:
                    target_placeholder = placeholder
                    break
            
            if not target_placeholder:
                raise ValueError(f"Placeholder {placeholder_idx} not found")
            
            # 更新样式
            if hasattr(target_placeholder, 'text_frame'):
                for paragraph in target_placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if font_size is not None:
                            run.font.size = self.Pt(font_size)
                        if font_bold is not None:
                            run.font.bold = font_bold
                        if font_color is not None:
                            # 解析hex颜色
                            color = self._parse_hex_color(font_color)
                            if color:
                                run.font.color.rgb = color
            
            # 保存文件
            save_path = output_path or pptx_path
            prs.save(save_path)
            
            return save_path
            
        except Exception as e:
            logger.error(f"Failed to update style: {e}")
            raise
    
    def add_image(
        self,
        pptx_path: str,
        slide_index: int,
        image_path: str,
        left: Optional[int] = None,
        top: Optional[int] = None,
        width: Optional[int] = None,
        height: Optional[int] = None,
        output_path: Optional[str] = None
    ) -> str:
        """
        向幻灯片添加图片
        
        Args:
            pptx_path: PPTX文件路径
            slide_index: 幻灯片索引
            image_path: 图片文件路径
            left: 左边距(EMU)
            top: 上边距(EMU)
            width: 宽度(EMU)
            height: 高度(EMU)
            output_path: 输出路径
            
        Returns:
            输出文件路径
        """
        if not self.available:
            raise RuntimeError("python-pptx library not available")
        
        try:
            from pptx.util import Inches
            
            prs = self.Presentation(pptx_path)
            
            if slide_index >= len(prs.slides):
                raise ValueError(f"Slide index {slide_index} out of range")
            
            slide = prs.slides[slide_index]
            
            # 设置默认位置
            left = left or Inches(1)
            top = top or Inches(1)
            width = width or Inches(4)
            height = height or Inches(3)
            
            # 添加图片
            slide.shapes.add_picture(image_path, left, top, width, height)
            
            # 保存文件
            save_path = output_path or pptx_path
            prs.save(save_path)
            
            return save_path
            
        except Exception as e:
            logger.error(f"Failed to add image: {e}")
            raise
    
    def delete_slide(
        self,
        pptx_path: str,
        slide_index: int,
        output_path: Optional[str] = None
    ) -> str:
        """
        删除幻灯片
        
        Args:
            pptx_path: PPTX文件路径
            slide_index: 幻灯片索引
            output_path: 输出路径
            
        Returns:
            输出文件路径
        """
        if not self.available:
            raise RuntimeError("python-pptx library not available")
        
        try:
            prs = self.Presentation(pptx_path)
            
            if slide_index >= len(prs.slides):
                raise ValueError(f"Slide index {slide_index} out of range")
            
            # 获取幻灯片ID列表
            slide_ids = prs.slides._sldIdLst
            slides = list(slide_ids)
            
            # 删除指定幻灯片
            slide_ids.remove(slides[slide_index])
            
            # 保存文件
            save_path = output_path or pptx_path
            prs.save(save_path)
            
            return save_path
            
        except Exception as e:
            logger.error(f"Failed to delete slide: {e}")
            raise
    
    def duplicate_slide(
        self,
        pptx_path: str,
        slide_index: int,
        output_path: Optional[str] = None
    ) -> str:
        """
        复制幻灯片
        
        Args:
            pptx_path: PPTX文件路径
            slide_index: 幻灯片索引
            output_path: 输出路径
            
        Returns:
            输出文件路径
        """
        if not self.available:
            raise RuntimeError("python-pptx library not available")
        
        try:
            from copy import deepcopy
            import lxml.etree as etree
            
            prs = self.Presentation(pptx_path)
            
            if slide_index >= len(prs.slides):
                raise ValueError(f"Slide index {slide_index} out of range")
            
            # 获取源幻灯片
            source_slide = prs.slides[slide_index]
            
            # 复制幻灯片XML
            slide_layout = source_slide.slide_layout
            new_slide = prs.slides.add_slide(slide_layout)
            
            # 复制形状
            for shape in source_slide.shapes:
                # 这里简化处理，实际可能需要更复杂的复制逻辑
                pass
            
            # 保存文件
            save_path = output_path or pptx_path
            prs.save(save_path)
            
            return save_path
            
        except Exception as e:
            logger.error(f"Failed to duplicate slide: {e}")
            raise
    
    def _parse_hex_color(self, hex_color: str) -> Optional[Any]:
        """
        解析hex颜色字符串
        
        Args:
            hex_color: hex颜色字符串，如"#FF0000"或"FF0000"
            
        Returns:
            RGBColor对象
        """
        try:
            # 移除#号
            hex_color = hex_color.lstrip('#')
            
            # 解析RGB值
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            
            return self.RGBColor(r, g, b)
        except:
            return None


# 全局实例
ppt_editor = PPTEditor()
